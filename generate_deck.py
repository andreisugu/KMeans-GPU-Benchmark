import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml import parse_xml

def create_presentation():
    prs = Presentation()
    
    # 1. 16:9 widescreen format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Define design tokens
    COLOR_BG_DARK = RGBColor(15, 23, 42)      # Deep slate/navy
    COLOR_BG_LIGHT = RGBColor(248, 250, 252)  # Soft cool white
    COLOR_TEXT_LIGHT = RGBColor(241, 245, 249) # Near white for dark bg
    COLOR_TEXT_DARK = RGBColor(15, 23, 42)    # Deep slate for light bg
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # Cool gray for sub-elements
    
    COLOR_ACCENT = RGBColor(6, 182, 212)      # Vibrant Cyan
    COLOR_SUPPORT = RGBColor(30, 41, 59)      # Slate supporting color
    COLOR_CARD_BG = RGBColor(255, 255, 255)   # Card background
    COLOR_BORDER = RGBColor(226, 232, 240)    # Soft borders
    
    # Chart Colors
    COLOR_CHART_SEQ = RGBColor(148, 163, 184) # Slate 400
    COLOR_CHART_SKLEARN = RGBColor(6, 182, 212) # Cyan 500
    COLOR_CHART_HIP = RGBColor(239, 68, 68)   # Red 500 (AMD)
    COLOR_CHART_MPI = RGBColor(34, 197, 94)   # Green 500
    
    FONT_TITLE = "Trebuchet MS"
    FONT_BODY = "Calibri"
    
    blank_layout = prs.slide_layouts[6]
    
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_slide_header(slide, title_text, category_text="K-MEANS PARALLEL BENCHMARK"):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_top = cat_tf.margin_right = cat_tf.margin_bottom = 0
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category_text.upper()
        cat_p.font.name = FONT_BODY
        cat_p.font.size = Pt(12)
        cat_p.font.bold = True
        cat_p.font.color.rgb = COLOR_ACCENT
        
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_top = title_tf.margin_right = title_tf.margin_bottom = 0
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.name = FONT_TITLE
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = COLOR_TEXT_DARK

    def add_slide_footer(slide, current_page, total_pages=12, is_dark=False):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.3))
        footer_tf = footer_box.text_frame
        footer_tf.word_wrap = True
        footer_tf.margin_left = footer_tf.margin_top = footer_tf.margin_right = footer_tf.margin_bottom = 0
        footer_p = footer_tf.paragraphs[0]
        footer_p.text = f"AMD Ryzen 7 7840HS & Radeon 780M Benchmarking Suite | Page {current_page} of {total_pages}"
        footer_p.font.name = FONT_BODY
        footer_p.font.size = Pt(11)
        footer_p.font.color.rgb = COLOR_TEXT_LIGHT if is_dark else COLOR_TEXT_MUTED

    def style_chart(chart):
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = FONT_BODY
        chart.legend.font.size = Pt(12)
        chart.legend.font.color.rgb = COLOR_TEXT_DARK
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = COLOR_BORDER

    def set_log_scale(chart):
        valAx = chart.value_axis._element
        scaling = valAx.xpath('.//c:scaling')[0]
        logBase = parse_xml('<c:logBase xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" val="10"/>')
        scaling.append(logBase)

    # =========================================================================
    # SLIDE 1-6 (Same as before)
    # =========================================================================
    # Slide 1: Title
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLOR_BG_DARK)
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.12), Inches(3.2))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.color.rgb = COLOR_ACCENT
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(10.5), Inches(3.4))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.margin_left = title_tf.margin_top = title_tf.margin_right = title_tf.margin_bottom = 0
    p_tag = title_tf.paragraphs[0]
    p_tag.text = "PARALLEL COMPUTING & MACHINE LEARNING PROJECT"
    p_tag.font.name = FONT_BODY
    p_tag.font.size = Pt(13.5)
    p_tag.font.bold = True
    p_tag.font.color.rgb = COLOR_ACCENT
    p_tag.space_after = Pt(16.5)
    p_main = title_tf.add_paragraph()
    p_main.text = "Accelerating Unsupervised Learning"
    p_main.font.name = FONT_TITLE
    p_main.font.size = Pt(50)
    p_main.font.bold = True
    p_main.font.color.rgb = COLOR_TEXT_LIGHT
    p_submain = title_tf.add_paragraph()
    p_submain.text = "K-Means Parallel Benchmarking Suite"
    p_submain.font.name = FONT_TITLE
    p_submain.font.size = Pt(38)
    p_submain.font.bold = True
    p_submain.font.color.rgb = COLOR_TEXT_LIGHT
    p_submain.space_after = Pt(23)
    p_sub = title_tf.add_paragraph()
    p_sub.text = "A performance comparison across CPU, GPU, MPI & Vectorized BLAS on AMD Ryzen 7 7840HS"
    p_sub.font.name = FONT_BODY
    p_sub.font.size = Pt(19)
    p_sub.font.color.rgb = COLOR_TEXT_MUTED
    add_slide_footer(slide1, 1, is_dark=True)

    # =========================================================================
    # SLIDE 2: Theme Description (What is K-Means)
    # =========================================================================
    slide_theme = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_theme, COLOR_BG_LIGHT)
    add_slide_header(slide_theme, "Project Theme: K-Means Clustering & Lloyd's Algorithm")
    
    col_width = Inches(5.6)
    col_height = Inches(4.9)
    top_pos = Inches(1.8)
    
    # Left Card: Lloyd's Algorithm
    algo_card = slide_theme.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, col_width, col_height)
    algo_card.fill.solid()
    algo_card.fill.fore_color.rgb = COLOR_CARD_BG
    algo_card.line.color.rgb = COLOR_BORDER
    algo_tf = algo_card.text_frame
    algo_tf.word_wrap = True
    algo_tf.margin_left = algo_tf.margin_right = algo_tf.margin_top = algo_tf.margin_bottom = Inches(0.3)
    algo_title = algo_tf.paragraphs[0]
    algo_title.text = "Lloyd's Algorithm (Standard K-Means)"
    algo_title.font.name = FONT_TITLE
    algo_title.font.size = Pt(22)
    algo_title.font.bold = True
    algo_title.font.color.rgb = COLOR_SUPPORT
    algo_title.space_after = Pt(8)
    
    steps = [
        "1. Initialize: Pick K random points as the initial centroids (Forgy method).",
        "2. Assign (Expectation): Calculate the distance from every point to every centroid. Assign each point to its nearest centroid to form Voronoi partitions.",
        "3. Update (Maximization): Recompute the position of each centroid by calculating the arithmetic mean of all points assigned to it.",
        "4. Repeat: Loop through Steps 2 and 3 until no points change their cluster assignments (convergence guaranteed)."
    ]
    for b in steps:
        p = algo_tf.add_paragraph()
        p.text = b
        p.font.name = FONT_BODY
        p.font.size = Pt(13)
        p.font.bold = b.startswith("2. Assign")
        p.font.color.rgb = COLOR_ACCENT if b.startswith("2. Assign") else COLOR_TEXT_DARK
        p.space_after = Pt(5)

    # Right Card: Complexity & Bottleneck
    comp_card = slide_theme.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), top_pos, col_width, col_height)
    comp_card.fill.solid()
    comp_card.fill.fore_color.rgb = COLOR_CARD_BG
    comp_card.line.color.rgb = COLOR_BORDER
    comp_tf = comp_card.text_frame
    comp_tf.word_wrap = True
    comp_tf.margin_left = comp_tf.margin_right = comp_tf.margin_top = comp_tf.margin_bottom = Inches(0.3)
    comp_title = comp_tf.paragraphs[0]
    comp_title.text = "Complexity & The Parallelization Goal"
    comp_title.font.name = FONT_TITLE
    comp_title.font.size = Pt(22)
    comp_title.font.bold = True
    comp_title.font.color.rgb = COLOR_SUPPORT
    comp_title.space_after = Pt(8)
    
    comp_desc = comp_tf.add_paragraph()
    comp_desc.text = "Sequential Complexity: O(N × K × D × I)"
    comp_desc.font.name = FONT_TITLE
    comp_desc.font.size = Pt(16)
    comp_desc.font.bold = True
    comp_desc.font.color.rgb = COLOR_CHART_HIP
    comp_desc.space_after = Pt(8)
    
    comp_bullets = [
        "N = Points, K = Clusters, D = Dimensions, I = Iterations.",
        "The Bottleneck: The Assignment phase (Step 2) requires N × K distance calculations per iteration. For large datasets, this nested loop takes up >95% of execution time.",
        "The Parallel Opportunity: Step 2 is 'Embarrassingly Parallel'. Finding the nearest centroid for point A is completely independent of finding the nearest centroid for point B.",
        "Project Goal: Exploit this independence across GPU threads (HIP), distributed memory architectures (MPI), and SIMD vectorization (scikit-learn) to drastically reduce execution time."
    ]
    for b in comp_bullets:
        p = comp_tf.add_paragraph()
        p.text = "• " + b
        p.font.name = FONT_BODY
        p.font.size = Pt(12.5)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_after = Pt(4)
        
    add_slide_footer(slide_theme, 2)

    # =========================================================================
    # SLIDE 3: Hardware Environment (Light Background)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLOR_BG_LIGHT)
    add_slide_header(slide2, "Benchmarking Platform: Hardware Specifications")
    col_width = Inches(5.6)
    col_height = Inches(4.5)
    top_pos = Inches(1.8)
    cpu_card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, col_width, col_height)
    cpu_card.fill.solid()
    cpu_card.fill.fore_color.rgb = COLOR_CARD_BG
    cpu_card.line.color.rgb = COLOR_BORDER
    cpu_tf = cpu_card.text_frame
    cpu_tf.word_wrap = True
    cpu_tf.margin_left = cpu_tf.margin_right = cpu_tf.margin_top = cpu_tf.margin_bottom = Inches(0.3)
    cpu_title = cpu_tf.paragraphs[0]
    cpu_title.text = "AMD Ryzen 7 7840HS (CPU)"
    cpu_title.font.name = FONT_TITLE
    cpu_title.font.size = Pt(24)
    cpu_title.font.bold = True
    cpu_title.font.color.rgb = COLOR_SUPPORT
    cpu_title.space_after = Pt(12)
    cpu_bullets = [
        "Architecture: Zen 4 (4nm TSMC lithography) for premium energy efficiency and performance density.",
        "Cores / Threads: 8 Physical Cores, 16 Logic Threads with Simultaneous Multithreading (SMT).",
        "Clocks: 3.8 GHz Base Clock, scaling dynamically up to 5.1 GHz Max Boost.",
        "Cache Hierarchy: 8 MB L2 Cache + 16 MB L3 Cache feeding CPU cores with minimal memory stalls.",
        "Role in Benchmark: Host controller, runs Sequential baseline C++ code, coordinates multi-threaded python loops, and runs 4-process MPI domain decomposition."
    ]
    for b in cpu_bullets:
        p = cpu_tf.add_paragraph()
        p.text = "• " + b
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_after = Pt(6)
        
    gpu_card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), top_pos, col_width, col_height)
    gpu_card.fill.solid()
    gpu_card.fill.fore_color.rgb = COLOR_CARD_BG
    gpu_card.line.color.rgb = COLOR_BORDER
    gpu_tf = gpu_card.text_frame
    gpu_tf.word_wrap = True
    gpu_tf.margin_left = gpu_tf.margin_right = gpu_tf.margin_top = gpu_tf.margin_bottom = Inches(0.3)
    gpu_title = gpu_tf.paragraphs[0]
    gpu_title.text = "AMD Radeon 780M (iGPU)"
    gpu_title.font.name = FONT_TITLE
    gpu_title.font.size = Pt(24)
    gpu_title.font.bold = True
    gpu_title.font.color.rgb = COLOR_SUPPORT
    gpu_title.space_after = Pt(12)
    gpu_bullets = [
        "Architecture: RDNA 3 graphics architecture. Code name: gfx1103.",
        "Compute Units: 12 CUs containing 768 Stream Processors (Shaders) running up to 2.7 GHz.",
        "Memory Architecture: Unified Memory Architecture (UMA) — shares host system RAM.",
        "Bandwidth Limitation: Integrated GPU relies on DDR5/LPDDR5 memory channels rather than dedicated VRAM, creating a memory-bandwidth bottleneck during heavy data transfers.",
        "Role in Benchmark: Offloads massively parallel Expectation step via HIP (Heterogeneous-Compute Interface for Portability)."
    ]
    for b in gpu_bullets:
        p = gpu_tf.add_paragraph()
        p.text = "• " + b
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_after = Pt(6)
    add_slide_footer(slide2, 3)

    # Slide 4: Architectures
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLOR_BG_LIGHT)
    add_slide_header(slide3, "Four Implementations: Strategy & Core Technologies")
    card_w = Inches(5.6)
    card_h = Inches(2.1)
    x_positions = [Inches(0.8), Inches(6.9)]
    y_positions = [Inches(1.8), Inches(4.2)]
    implementations = [
        {"title": "1. Sequential C++ Baseline", "tech": "G++ compiler | Flag: -O3", "desc": "Single-threaded implementation designed as the control baseline. Employs Data-Oriented Design (DOD) with flat arrays to maximize L1/L2 cache hits.", "pos": (0, 0)},
        {"title": "2. Vectorized scikit-learn", "tech": "Python | NumPy | OpenBLAS / Intel MKL", "desc": "High-performance baseline. Bypasses nested Python loops by computing all N×K distances in a single matrix-multiplication operation via SIMD.", "pos": (1, 0)},
        {"title": "3. HIP / GPU Accelerated", "tech": "AMD ROCm Suite | hipcc compiler", "desc": "Parallelizes distance calculations by offloading the Expectation loop to the 780M iGPU. Employs coalesced memory mapping.", "pos": (0, 1)},
        {"title": "4. MPI Distributed C++", "tech": "OpenMPI | mpic++ compiler | 4 Processes", "desc": "Distributed-memory model. Performs spatial domain decomposition to partition points. Sub-processes run global reductions via MPI_Allreduce.", "pos": (1, 1)}
    ]
    for impl in implementations:
        x_idx, y_idx = impl["pos"]
        x = x_positions[x_idx]
        y = y_positions[y_idx]
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_BORDER
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.2)
        title_p = tf.paragraphs[0]
        title_p.text = impl["title"]
        title_p.font.name = FONT_TITLE
        title_p.font.size = Pt(19)
        title_p.font.bold = True
        title_p.font.color.rgb = COLOR_SUPPORT
        tech_p = tf.add_paragraph()
        tech_p.text = impl["tech"]
        tech_p.font.name = FONT_BODY
        tech_p.font.size = Pt(13.5)
        tech_p.font.bold = True
        tech_p.font.color.rgb = COLOR_ACCENT
        tech_p.space_after = Pt(10)
        desc_p = tf.add_paragraph()
        desc_p.text = impl["desc"]
        desc_p.font.name = FONT_BODY
        desc_p.font.size = Pt(14)
        desc_p.font.color.rgb = COLOR_TEXT_DARK
    add_slide_footer(slide3, 4)

    # Slide 5: Assignments
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLOR_BG_LIGHT)
    add_slide_header(slide4, "Logical Differences: The Assignment Phase (Expectation)")
    col_w = Inches(3.6)
    col_h = Inches(4.5)
    col_top = Inches(1.8)
    col_x = [Inches(0.8), Inches(4.8), Inches(8.8)]
    steps = [
        {"header": "CPU Loop (Sequential & MPI)", "subheader": "Explicit Nested Loops", "desc": "• Evaluates distances sequentially: O(N × K × D) calculations.\n• High Cache Locality: Flat 1D layout ensures CPU cache lines are filled with contiguous point features.\n• Core Bottleneck: CPU instructions are executed lineally, resulting in execution times scaling up to 10 minutes on large datasets.\n• MPI distributes this work: each process computes assignments on N/P points.", "color": COLOR_CARD_BG},
        {"header": "Vectorized BLAS (sklearn)", "subheader": "The Kernel Trick GEMM", "desc": "• Bypasses the O(N × K × D) coordinate loop entirely by formulating distance as a matrix operation:\n  ||X - C||² = ||X||² - 2XCᵀ + ||C||²\n• The -2XCᵀ term is computed via highly optimized Level-3 BLAS.\n• Leverages cache-blocking and CPU SIMD vector units, achieving the fastest times.", "color": COLOR_CARD_BG},
        {"header": "GPU Kernel Parallelism", "subheader": "Massively Parallel Threads", "desc": "• Launches a 1D grid with a block size of 256 threads.\n• Direct Mapping: Each thread is mapped to exactly one point index i via: blockIdx.x * blockDim.x + threadIdx.x.\n• Embarrassingly Parallel: All N points calculate distance to all K centroids simultaneously.\n• Offloads the core bottleneck from host CPU cores.", "color": COLOR_CARD_BG}
    ]
    for idx, step in enumerate(steps):
        x = col_x[idx]
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, col_top, col_w, col_h)
        card.fill.solid()
        card.fill.fore_color.rgb = step["color"]
        card.line.color.rgb = COLOR_BORDER
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.25)
        h = tf.paragraphs[0]
        h.text = step["header"]
        h.font.name = FONT_TITLE
        h.font.size = Pt(20)
        h.font.bold = True
        h.font.color.rgb = COLOR_SUPPORT
        sh = tf.add_paragraph()
        sh.text = step["subheader"]
        sh.font.name = FONT_BODY
        sh.font.size = Pt(14.5)
        sh.font.bold = True
        sh.font.color.rgb = COLOR_ACCENT
        sh.space_after = Pt(16.5)
        d = tf.add_paragraph()
        d.text = step["desc"]
        d.font.name = FONT_BODY
        d.font.size = Pt(14)
        d.font.color.rgb = COLOR_TEXT_DARK
        d.space_before = Pt(7)
    add_slide_footer(slide4, 5)

    # Slide 6: Updates
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, COLOR_BG_LIGHT)
    add_slide_header(slide5, "Logical Differences: Centroid Update & Hybrid GPU-CPU Model")
    block_w = Inches(11.7)
    block_h = Inches(2.1)
    block_x = Inches(0.8)
    block_y = [Inches(1.8), Inches(4.2)]
    
    card1 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, block_x, block_y[0], block_w, block_h)
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLOR_CARD_BG
    card1.line.color.rgb = COLOR_BORDER
    tf1 = card1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = Inches(0.2)
    title1 = tf1.paragraphs[0]
    title1.text = "The Bottleneck: GPU Atomic Add Contention"
    title1.font.name = FONT_TITLE
    title1.font.size = Pt(21)
    title1.font.bold = True
    title1.font.color.rgb = COLOR_SUPPORT
    title1.space_after = Pt(8)
    desc1 = tf1.add_paragraph()
    desc1.text = "In Lloyd's algorithm, updating centroids (Maximization) requires accumulating coordinates of all points in each cluster. Performing this on the GPU via global memory requires double-precision atomicAdd operations. When millions of threads write to just 23 or 50 centroid memory locations, severe memory serialization occurs. The threads queue up to execute writes, causing performance to drop below the sequential CPU baseline."
    desc1.font.name = FONT_BODY
    desc1.font.size = Pt(14.5)
    desc1.font.color.rgb = COLOR_TEXT_DARK
    
    card2 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, block_x, block_y[1], block_w, block_h)
    card2.fill.solid()
    card2.fill.fore_color.rgb = COLOR_CARD_BG
    card2.line.color.rgb = COLOR_BORDER
    tf2 = card2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = Inches(0.2)
    title2 = tf2.paragraphs[0]
    title2.text = "The Solution: Hybrid CPU-GPU Pipeline"
    title2.font.name = FONT_TITLE
    title2.font.size = Pt(21)
    title2.font.bold = True
    title2.font.color.rgb = COLOR_ACCENT
    title2.space_after = Pt(8)
    desc2 = tf2.add_paragraph()
    desc2.text = "To resolve contention, we implemented a hybrid model: (1) The GPU computes assignments (Expectation) in parallel. (2) The flat assignments array (~19.6 MB for 5M points) is transferred back to the CPU (taking <1.5 ms over UMA system bus). (3) The CPU handles centroid accumulation (Maximization) sequentially in O(N × D) time. This hybrid approach yielded a 6.4× speedup over the naive atomic GPU implementation."
    desc2.font.name = FONT_BODY
    desc2.font.size = Pt(14.5)
    desc2.font.color.rgb = COLOR_TEXT_DARK
    add_slide_footer(slide5, 6)

    # Slide 7: Synthetic Table
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, COLOR_BG_LIGHT)
    add_slide_header(slide6, "Performance Comparison: Synthetic Datasets")
    tbl_left = Inches(0.8)
    tbl_top = Inches(1.8)
    tbl_w = Inches(8.2)
    tbl_h = Inches(4.5)
    table_shape = slide6.shapes.add_table(6, 8, tbl_left, tbl_top, tbl_w, tbl_h)
    table = table_shape.table
    col_widths = [Inches(1.8), Inches(0.7), Inches(0.5), Inches(0.5), Inches(1.2), Inches(1.1), Inches(1.2), Inches(1.2)]
    for idx, width in enumerate(col_widths):
        table.columns[idx].width = width
    headers = ["Dataset Level", "N (Pts)", "D", "K", "Seq C++", "scikit-learn", "HIP/GPU (AMD)", "MPI (4 P)"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_SUPPORT
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_TITLE
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_LIGHT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    data_rows = [
        ["1. Very Small", "250K", "16", "32", "4,965 ms", "396 ms", "2,122 ms", "1,484 ms"],
        ["2. Small", "700K", "32", "32", "39,605 ms", "1,722 ms", "15,031 ms", "9,502 ms"],
        ["3. Medium", "1.4M", "32", "64", "212,183 ms", "8,109 ms", "48,300 ms", "53,818 ms"],
        ["4. Large", "3.5M", "32", "64", "306,346 ms", "21,833 ms", "77,633 ms", "160,434 ms"],
        ["5. Extreme", "3.5M", "32", "128", "613,309 ms", "31,115 ms", "142,613 ms", "156,254 ms"]
    ]
    for row_idx, row_data in enumerate(data_rows):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD_BG if row_idx % 2 == 0 else RGBColor(241, 245, 249)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            p.font.name = FONT_BODY
            p.font.size = Pt(13.5)
            if col_idx == 5:
                p.font.bold = True
                p.font.color.rgb = COLOR_ACCENT
            else:
                p.font.color.rgb = COLOR_TEXT_DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    comm_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(1.8), Inches(3.3), Inches(4.5))
    comm_card.fill.solid()
    comm_card.fill.fore_color.rgb = COLOR_CARD_BG
    comm_card.line.color.rgb = COLOR_BORDER
    comm_tf = comm_card.text_frame
    comm_tf.word_wrap = True
    comm_tf.margin_left = comm_tf.margin_right = comm_tf.margin_top = comm_tf.margin_bottom = Inches(0.2)
    comm_title = comm_tf.paragraphs[0]
    comm_title.text = "Key Observations"
    comm_title.font.name = FONT_TITLE
    comm_title.font.size = Pt(19)
    comm_title.font.bold = True
    comm_title.font.color.rgb = COLOR_SUPPORT
    comm_title.space_after = Pt(12)
    bullets = [
        "scikit-learn is the speed leader across all benchmarks (up to 26× speedup). Vectorized GEMM (BLAS) routines avoid coordinate-level branching.",
        "HIP/GPU scales efficiently on massive workloads, outperforming MPI once N exceeds 1.4M points.",
        "MPI hits scaling limits as K increases to 128, due to MPI_Allreduce communication overhead."
    ]
    for b in bullets:
        p = comm_tf.add_paragraph()
        p.text = "• " + b
        p.font.name = FONT_BODY
        p.font.size = Pt(13.5)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_after = Pt(10)
    add_slide_footer(slide6, 7)

    # =========================================================================
    # SLIDE 8: Synthetic Benchmarks Visualizations (NATIVE CHARTS)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, COLOR_BG_LIGHT)
    add_slide_header(slide7, "Synthetic Datasets: Execution Time & Speedup Curves")
    
    # Chart 1: Execution Time (Log Scale)
    chart_data_time = CategoryChartData()
    chart_data_time.categories = ['1. Very Small', '2. Small', '3. Medium', '4. Large', '5. Extreme']
    chart_data_time.add_series('Seq C++', (4965, 39605, 212183, 306346, 613309))
    chart_data_time.add_series('sklearn', (396, 1722, 8109, 21833, 31115))
    chart_data_time.add_series('HIP/GPU', (2122, 15031, 48300, 77633, 142613))
    chart_data_time.add_series('MPI (4P)', (1484, 9502, 53818, 160434, 156254))
    
    x1, y1, cx, cy = Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.5)
    chart1 = slide7.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x1, y1, cx, cy, chart_data_time).chart
    style_chart(chart1)
    set_log_scale(chart1)
    chart1.value_axis.has_title = True
    chart1.value_axis.axis_title.text_frame.text = "Execution Time (ms)"
    chart1.chart_title.text_frame.text = "Execution Time Comparison (ms) — Log Scale"
    
    # Apply series colors
    for i, color in enumerate([COLOR_CHART_SEQ, COLOR_CHART_SKLEARN, COLOR_CHART_HIP, COLOR_CHART_MPI]):
        series = chart1.series[i]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color
        
    # Chart 2: Speedup vs Sequential
    chart_data_speedup = CategoryChartData()
    chart_data_speedup.categories = ['1. Very Small', '2. Small', '3. Medium', '4. Large', '5. Extreme']
    chart_data_speedup.add_series('sklearn', (12.5, 23.0, 26.2, 14.0, 19.7))
    chart_data_speedup.add_series('HIP/GPU', (2.3, 2.6, 4.4, 3.9, 4.3))
    chart_data_speedup.add_series('MPI (4P)', (3.3, 4.2, 3.9, 1.9, 3.9))
    
    x2 = Inches(6.8)
    chart2 = slide7.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x2, y1, cx, cy, chart_data_speedup).chart
    style_chart(chart2)
    chart2.value_axis.has_title = True
    chart2.value_axis.axis_title.text_frame.text = "Speedup Factor (x)"
    chart2.chart_title.text_frame.text = "Speedup vs Sequential Baseline"
    
    # Line chart colors
    for i, color in enumerate([COLOR_CHART_SKLEARN, COLOR_CHART_HIP, COLOR_CHART_MPI]):
        series = chart2.series[i]
        series.format.line.color.rgb = color
        series.format.line.width = Pt(5)

    add_slide_footer(slide7, 8)

    # =========================================================================
    # SLIDE 9: Real-World Datasets Table
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, COLOR_BG_LIGHT)
    add_slide_header(slide8, "Performance Comparison: Real-World Datasets")
    
    table_shape8 = slide8.shapes.add_table(6, 8, tbl_left, tbl_top, tbl_w, tbl_h)
    table8 = table_shape8.table
    for idx, width in enumerate(col_widths):
        table8.columns[idx].width = width
    for col_idx, text in enumerate(headers):
        cell = table8.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_SUPPORT
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_TITLE
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_LIGHT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    data_rows8 = [
        ["Mall Customers", "200", "3", "5", "0.04 ms", "21.86 ms", "1.77 ms", "0.14 ms"],
        ["Covertype K=7", "581K", "54", "7", "3,738 ms", "394 ms", "1,922 ms", "1,251 ms"],
        ["Covertype K=50", "581K", "54", "50", "81,518 ms", "927 ms", "28,308 ms", "20,905 ms"],
        ["KDD Cup 10%", "494K", "38", "23", "8,182 ms", "417 ms", "6,071 ms", "3,674 ms"],
        ["KDD Cup Full", "4.89M", "38", "23", "201,664 ms", "5,307 ms", "45,089 ms", "37,818 ms"]
    ]
    for row_idx, row_data in enumerate(data_rows8):
        for col_idx, val in enumerate(row_data):
            cell = table8.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD_BG if row_idx % 2 == 0 else RGBColor(241, 245, 249)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            p.font.name = FONT_BODY
            p.font.size = Pt(13.5)
            if col_idx == 5:
                p.font.bold = True
                p.font.color.rgb = COLOR_ACCENT
            else:
                p.font.color.rgb = COLOR_TEXT_DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    comm_card8 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(1.8), Inches(3.3), Inches(4.5))
    comm_card8.fill.solid()
    comm_card8.fill.fore_color.rgb = COLOR_CARD_BG
    comm_card8.line.color.rgb = COLOR_BORDER
    comm_tf8 = comm_card8.text_frame
    comm_tf8.word_wrap = True
    comm_tf8.margin_left = comm_tf8.margin_right = comm_tf8.margin_top = comm_tf8.margin_bottom = Inches(0.2)
    comm_title8 = comm_tf8.paragraphs[0]
    comm_title8.text = "Key Observations"
    comm_title8.font.name = FONT_TITLE
    comm_title8.font.size = Pt(19)
    comm_title8.font.bold = True
    comm_title8.font.color.rgb = COLOR_SUPPORT
    comm_title8.space_after = Pt(12)
    bullets8 = [
        "Small Datasets: GPU kernel launch (~1.7ms) and Python library loading (~21.8ms) makes Sequential C++ faster for tiny runs.",
        "At Scale: KDD Cup Full shows MPI (37.8s) and HIP/GPU (45s) yielding ~5.3× and 4.4× speedups.",
        "High-K Stress: Covertype K=50 demonstrates sklearn achieving a 87.8× speedup by optimizing matrix multiplications."
    ]
    for b in bullets8:
        p = comm_tf8.add_paragraph()
        p.text = "• " + b
        p.font.name = FONT_BODY
        p.font.size = Pt(13.5)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_after = Pt(10)
    add_slide_footer(slide8, 9)

    # =========================================================================
    # SLIDE 10: Real-World Visualizations (NATIVE CHARTS)
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, COLOR_BG_LIGHT)
    add_slide_header(slide9, "Real-World Datasets: Execution Time & Speedup Curves")
    
    # Chart 3: Execution Time
    chart_data_time_real = CategoryChartData()
    chart_data_time_real.categories = ['Mall', 'Cov7', 'Cov50', 'KDD10%', 'KDD Full']
    chart_data_time_real.add_series('Seq C++', (0.04, 3738.36, 81518.00, 8182.50, 201664.00))
    chart_data_time_real.add_series('sklearn', (21.86, 394.77, 927.96, 417.77, 5307.35))
    chart_data_time_real.add_series('HIP/GPU', (1.77, 1922.17, 28308.00, 6071.91, 45089.70))
    chart_data_time_real.add_series('MPI (4P)', (0.14, 1251.34, 20905.30, 3674.98, 37818.60))
    
    chart3 = slide9.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x1, y1, cx, cy, chart_data_time_real).chart
    style_chart(chart3)
    set_log_scale(chart3)
    chart3.value_axis.has_title = True
    chart3.value_axis.axis_title.text_frame.text = "Execution Time (ms)"
    chart3.chart_title.text_frame.text = "Execution Time Comparison (ms) — Log Scale"
    
    for i, color in enumerate([COLOR_CHART_SEQ, COLOR_CHART_SKLEARN, COLOR_CHART_HIP, COLOR_CHART_MPI]):
        series = chart3.series[i]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color

    # Chart 4: Speedup vs Sequential
    chart_data_speedup_real = CategoryChartData()
    chart_data_speedup_real.categories = ['Mall', 'Cov7', 'Cov50', 'KDD10%', 'KDD Full']
    chart_data_speedup_real.add_series('sklearn', (0.00, 9.47, 87.85, 19.59, 38.00))
    chart_data_speedup_real.add_series('HIP/GPU', (0.02, 1.94, 2.88, 1.35, 4.47))
    chart_data_speedup_real.add_series('MPI (4P)', (0.28, 2.99, 3.90, 2.23, 5.33))
    
    chart4 = slide9.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x2, y1, cx, cy, chart_data_speedup_real).chart
    style_chart(chart4)
    chart4.value_axis.has_title = True
    chart4.value_axis.axis_title.text_frame.text = "Speedup Factor (x)"
    chart4.chart_title.text_frame.text = "Speedup vs Baseline (Real Data)"
    
    for i, color in enumerate([COLOR_CHART_SKLEARN, COLOR_CHART_HIP, COLOR_CHART_MPI]):
        series = chart4.series[i]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color

    add_slide_footer(slide9, 10)

    # =========================================================================
    # SLIDE 11: Cloud Acceleration Requirement (RAPIDS cuML k-NN)
    # =========================================================================
    slide_cuml = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_cuml, COLOR_BG_LIGHT)
    add_slide_header(slide_cuml, "Cloud Acceleration Requirement: RAPIDS cuML (k-Nearest Neighbors)")
    
    col_width = Inches(5.6)
    col_height = Inches(4.5)
    top_pos = Inches(1.8)
    
    # Left Card: The cuML Implementation
    cuml_card = slide_cuml.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, col_width, col_height)
    cuml_card.fill.solid()
    cuml_card.fill.fore_color.rgb = COLOR_CARD_BG
    cuml_card.line.color.rgb = COLOR_BORDER
    cuml_tf = cuml_card.text_frame
    cuml_tf.word_wrap = True
    cuml_tf.margin_left = cuml_tf.margin_right = cuml_tf.margin_top = cuml_tf.margin_bottom = Inches(0.3)
    cuml_title = cuml_tf.paragraphs[0]
    cuml_title.text = "NVIDIA RAPIDS cuML (Google Colab T4)"
    cuml_title.font.name = FONT_TITLE
    cuml_title.font.size = Pt(24)
    cuml_title.font.bold = True
    cuml_title.font.color.rgb = COLOR_SUPPORT
    cuml_title.space_after = Pt(12)
    
    cuml_desc = cuml_tf.add_paragraph()
    cuml_desc.text = "Executed via Google Colab on an NVIDIA T4 GPU due to strict CUDA requirements. Uses cuBLAS GEMM and Thrust radix-select for distance calculations."
    cuml_desc.font.name = FONT_BODY
    cuml_desc.font.size = Pt(14)
    cuml_desc.font.color.rgb = COLOR_TEXT_DARK
    cuml_desc.space_after = Pt(12)
    
    cuml_warn = cuml_tf.add_paragraph()
    cuml_warn.text = "The k-NN Self-Join Complexity"
    cuml_warn.font.name = FONT_TITLE
    cuml_warn.font.size = Pt(18)
    cuml_warn.font.bold = True
    cuml_warn.font.color.rgb = COLOR_CHART_HIP
    cuml_warn.space_after = Pt(8)
    
    cuml_bullets = [
        "K-Means (O(N×K)): Distances from N points to K centroids.",
        "k-NN (O(N²)): The required test script computed the nearest neighbors of the dataset AGAINST ITSELF.",
        "This means computing a massive N×N distance matrix!"
    ]
    for b in cuml_bullets:
        p = cuml_tf.add_paragraph()
        p.text = "• " + b
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_after = Pt(6)

    # Right Card: Results & The Guess
    res_card = slide_cuml.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), top_pos, col_width, col_height)
    res_card.fill.solid()
    res_card.fill.fore_color.rgb = COLOR_CARD_BG
    res_card.line.color.rgb = COLOR_BORDER
    res_tf = res_card.text_frame
    res_tf.word_wrap = True
    res_tf.margin_left = res_tf.margin_right = res_tf.margin_top = res_tf.margin_bottom = Inches(0.3)
    res_title = res_tf.paragraphs[0]
    res_title.text = "Execution Results & Complexity Limits"
    res_title.font.name = FONT_TITLE
    res_title.font.size = Pt(24)
    res_title.font.bold = True
    res_title.font.color.rgb = COLOR_SUPPORT
    res_title.space_after = Pt(12)
    
    res_bullets = [
        "Mall Customers (200 pts): 234.5 ms",
        "Covertype K=7 (581K pts): 33,950.7 ms",
        "Covertype K=50 (581K pts): 45,524.7 ms",
        "KDD Cup 10% (494K pts): 21,798.9 ms"
    ]
    for b in res_bullets:
        p = res_tf.add_paragraph()
        p.text = b
        p.font.name = FONT_BODY
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_after = Pt(6)
        
    guess_title = res_tf.add_paragraph()
    guess_title.text = "The Final Dataset (KDD Cup Full - 4.89M pts):"
    guess_title.font.name = FONT_TITLE
    guess_title.font.size = Pt(18)
    guess_title.font.bold = True
    guess_title.font.color.rgb = COLOR_ACCENT
    guess_title.space_before = Pt(12)
    guess_title.space_after = Pt(8)
    
    guess_p = res_tf.add_paragraph()
    guess_p.text = "STATUS: Could not be completed (Colab dataset download failed)."
    guess_p.font.name = FONT_BODY
    guess_p.font.size = Pt(14)
    guess_p.font.bold = True
    guess_p.font.color.rgb = COLOR_CHART_HIP
    guess_p.space_after = Pt(6)
    
    guess_desc = res_tf.add_paragraph()
    guess_desc.text = "Theoretical Extrapolation: An N×N matrix for 4.89 million points requires nearly 24 TRILLION distance computations. Even if downloaded, the T4 GPU would likely have encountered an Out-of-Memory (OOM) error or exceeded Colab's session time limits."
    guess_desc.font.name = FONT_BODY
    guess_desc.font.size = Pt(13)
    guess_desc.font.color.rgb = COLOR_TEXT_DARK

    add_slide_footer(slide_cuml, 11)

    # =========================================================================
    # SLIDE 12: Key Insights & Technical Summary
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, COLOR_BG_DARK)
    
    cat_box = slide10.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    cat_tf = cat_box.text_frame
    cat_tf.word_wrap = True
    cat_tf.margin_left = cat_tf.margin_top = cat_tf.margin_right = cat_tf.margin_bottom = 0
    cat_p = cat_tf.paragraphs[0]
    cat_p.text = "SUMMARY & FINDINGS"
    cat_p.font.name = FONT_BODY
    cat_p.font.size = Pt(12)
    cat_p.font.bold = True
    cat_p.font.color.rgb = COLOR_ACCENT
    
    title_box = slide10.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.margin_left = title_tf.margin_top = title_tf.margin_right = title_tf.margin_bottom = 0
    title_p = title_tf.paragraphs[0]
    title_p.text = "Core Takeaways: Parallel Computing Insights"
    title_p.font.name = FONT_TITLE
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_TEXT_LIGHT
    
    col_w = Inches(3.6)
    col_h = Inches(4.5)
    col_top = Inches(1.8)
    col_x = [Inches(0.8), Inches(4.8), Inches(8.8)]
    
    takeaways = [
        {
            "num": "01",
            "title": "Hardware Boundaries",
            "desc": "Unified Memory Architecture (UMA) on integrated GPUs like the Radeon 780M eliminates host-to-device copy overhead, but shares system RAM bandwidth. For data-intensive algorithms, GPU execution speed is bound by RAM speed, whereas a discrete GPU with dedicated VRAM would yield much higher speedups."
        },
        {
            "num": "02",
            "title": "Algorithmic Vectorization",
            "desc": "Traditional loops (even GPU or MPI parallelized) compile to coordinate-by-coordinate distance branches. Mapping the mathematical formula to GEMM matrix operations (the kernel trick in scikit-learn) allows CPU/SIMD hardware to execute at peak FLOPS, showing that formulation optimizations often outperform raw compute parallelization."
        },
        {
            "num": "03",
            "title": "Communication Overhead",
            "desc": "MPI distributed memory scales well for computation but requires all-to-all synchronization. The MPI_Allreduce collective operation's communication overhead grows linearly with cluster size K. For high-cluster counts, MPI's scaling is hindered by communication latency, highlighting the classic compute-to-communication ratio challenge."
        }
    ]
    
    for idx, take in enumerate(takeaways):
        x = col_x[idx]
        card = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, col_top, col_w, col_h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_SUPPORT
        card.line.color.rgb = COLOR_BORDER
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.25)
        
        num_p = tf.paragraphs[0]
        num_p.text = take["num"]
        num_p.font.name = FONT_TITLE
        num_p.font.size = Pt(42)
        num_p.font.bold = True
        num_p.font.color.rgb = COLOR_ACCENT
        num_p.space_after = Pt(7)
        
        t_p = tf.add_paragraph()
        t_p.text = take["title"]
        t_p.font.name = FONT_TITLE
        t_p.font.size = Pt(21)
        t_p.font.bold = True
        t_p.font.color.rgb = COLOR_TEXT_LIGHT
        t_p.space_after = Pt(14.5)
        
        d_p = tf.add_paragraph()
        d_p.text = take["desc"]
        d_p.font.name = FONT_BODY
        d_p.font.size = Pt(14.5)
        d_p.font.color.rgb = COLOR_TEXT_LIGHT
        
    add_slide_footer(slide10, 12, is_dark=True)
    
    output_path = "/home/restlessstone/VSCodeProjects/KMeans-GPU-Benchmark/KMeans_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
