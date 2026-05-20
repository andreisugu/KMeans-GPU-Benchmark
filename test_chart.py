from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.oxml import parse_xml

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
chart_data = CategoryChartData()
chart_data.categories = ['A', 'B']
chart_data.add_series('S1', (10, 1000))
chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(4), Inches(3), chart_data).chart

try:
    valAx = chart.value_axis._element
    scaling = valAx.xpath('.//c:scaling')[0]
    logBase = parse_xml('<c:logBase xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" val="10"/>')
    scaling.append(logBase)
    print("Log base added via XML")
except Exception as e:
    print(f"Error: {e}")

prs.save('test.pptx')
